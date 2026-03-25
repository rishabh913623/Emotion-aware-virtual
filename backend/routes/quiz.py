"""Quiz generation endpoint based on emotion history."""
from collections import Counter
import io
import re

from flask import Blueprint, request, jsonify
from PyPDF2 import PdfReader
from pptx import Presentation

from utils.db import fetch_emotions_for_student

quiz_bp = Blueprint("quiz", __name__)

QUESTION_BANK = {
    "Easy": [
        {
            "question": "What is 2 + 2?",
            "options": ["3", "4", "5", "6"],
            "answer": "4",
        },
        {
            "question": "Which planet is known as the Red Planet?",
            "options": ["Earth", "Mars", "Jupiter", "Venus"],
            "answer": "Mars",
        },
    ],
    "Medium": [
        {
            "question": "What does HTTP stand for?",
            "options": ["HyperText Transfer Protocol", "HighText Transfer Protocol", "Hyperlink Transfer Process", "Host Transfer Protocol"],
            "answer": "HyperText Transfer Protocol",
        },
        {
            "question": "Which data structure uses FIFO order?",
            "options": ["Stack", "Queue", "Tree", "Graph"],
            "answer": "Queue",
        },
    ],
    "Hard": [
        {
            "question": "What is the time complexity of binary search?",
            "options": ["O(n)", "O(log n)", "O(n log n)", "O(1)"],
            "answer": "O(log n)",
        },
        {
            "question": "Which layer type is commonly used for image feature extraction?",
            "options": ["Dense", "Convolutional", "Dropout", "Embedding"],
            "answer": "Convolutional",
        },
    ],
}

EMOTION_TO_DIFFICULTY = {
    "Confused": "Easy",
    "Engaged": "Hard",
    "Neutral": "Medium",
    "Bored": "Easy",
    "Distracted": "Easy",
}


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text content from a PDF file."""
    reader = PdfReader(io.BytesIO(file_bytes))
    return "\n".join((page.extract_text() or "") for page in reader.pages).strip()


def extract_text_from_ppt(file_bytes: bytes) -> str:
    """Extract text content from a PPT/PPTX file."""
    presentation = Presentation(io.BytesIO(file_bytes))
    chunks = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks).strip()


def sentence_chunks(text: str) -> list[str]:
    """Return cleaned sentence-like chunks for quiz generation."""
    normalized = re.sub(r"\s+", " ", (text or "")).strip()
    if not normalized:
        return []
    chunks = re.split(r"(?<=[.!?])\s+", normalized)
    return [chunk.strip() for chunk in chunks if len(chunk.strip()) > 35]


def pick_options(answer: str, source_text: str) -> list[str]:
    """Create lightweight distractors from document vocabulary."""
    words = [
        word
        for word in re.findall(r"[A-Za-z][A-Za-z0-9_-]{3,}", source_text or "")
        if word.lower() != answer.lower()
    ]
    unique_words = []
    seen = set()
    for word in words:
        lowered = word.lower()
        if lowered not in seen:
            seen.add(lowered)
            unique_words.append(word)

    distractors = unique_words[:3]
    while len(distractors) < 3:
        distractors.append(f"Option {len(distractors) + 1}")

    options = [answer, *distractors[:3]]
    return options


def generate_mcq_from_text(text: str, max_questions: int = 5) -> list[dict]:
    """Generate basic MCQs from extracted file text."""
    chunks = sentence_chunks(text)
    questions = []

    for chunk in chunks[:max_questions]:
        keywords = [word for word in re.findall(r"[A-Za-z][A-Za-z0-9_-]{4,}", chunk)]
        if not keywords:
            continue

        answer = max(keywords, key=len)
        prompt = chunk.replace(answer, "_____", 1)
        options = pick_options(answer, text)

        questions.append(
            {
                "question": f"Fill in the blank: {prompt}",
                "options": options,
                "answer": answer,
            }
        )

    if not questions:
        return [
            {
                "question": "What is the main topic discussed in the uploaded material?",
                "options": ["Concept A", "Concept B", "Concept C", "Concept D"],
                "answer": "Concept A",
            }
        ]

    return questions


@quiz_bp.route("/generate-quiz", methods=["GET"])
def generate_quiz():
    """Generate quiz questions based on recent emotions."""
    try:
        student_id_raw = request.args.get("student_id", "1")
        try:
            student_id = int(student_id_raw)
        except (TypeError, ValueError):
            return jsonify({"error": "student_id must be an integer"}), 400

        history = fetch_emotions_for_student(student_id, limit=50)

        if history:
            recent = [item["emotion"] for item in history[:10]]
            most_recent_emotion = Counter(recent).most_common(1)[0][0]
        else:
            most_recent_emotion = "Neutral"

        difficulty = EMOTION_TO_DIFFICULTY.get(most_recent_emotion, "Medium")
        questions = QUESTION_BANK[difficulty]

        return jsonify({
            "student_id": student_id,
            "emotion": most_recent_emotion,
            "difficulty": difficulty,
            "questions": questions,
        })
    except Exception as exc:  # pragma: no cover - defensive error handling
        return jsonify({"error": "Failed to generate quiz", "details": str(exc)}), 500


@quiz_bp.route("/upload", methods=["POST"])
def upload_quiz_source():
    """Accept PDF/PPT uploads and generate MCQs from extracted text."""
    try:
        if "file" not in request.files:
            return jsonify({"error": "file is required"}), 400

        file = request.files["file"]
        filename = (file.filename or "").lower()
        file_bytes = file.read()

        if not filename:
            return jsonify({"error": "Invalid file name"}), 400
        if not file_bytes:
            return jsonify({"error": "Uploaded file is empty"}), 400

        if filename.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(file_bytes)
        elif filename.endswith(".ppt") or filename.endswith(".pptx"):
            extracted_text = extract_text_from_ppt(file_bytes)
        else:
            return jsonify({"error": "Only .pdf, .ppt, .pptx are supported"}), 400

        questions = generate_mcq_from_text(extracted_text)

        return jsonify({
            "source": file.filename,
            "questions": questions,
            "text_length": len(extracted_text),
        })
    except Exception as exc:  # pragma: no cover - defensive error handling
        return jsonify({"error": "Failed to process uploaded file", "details": str(exc)}), 500
